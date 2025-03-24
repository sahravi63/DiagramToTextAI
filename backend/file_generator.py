from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import os

output_dir = "output"
os.makedirs(output_dir, exist_ok=True)

def generate_pdf(text, filename="output.pdf"):
    path = os.path.join(output_dir, filename)
    c = canvas.Canvas(path, pagesize=letter)
    
    # Formatting adjustments
    x, y = 100, 750  # Start position
    max_width = 450  # Set max width for wrapping
    lines = text.split("\n")

    for line in lines:
        if y <= 50:  # New page if out of space
            c.showPage()
            y = 750

        c.drawString(x, y, line[:80])  # Limiting characters per line
        y -= 20  # Adjust line spacing

    c.save()
    return filename

def generate_docx(text, filename="output.docx"):
    path = os.path.join(output_dir, filename)
    doc = Document()
    doc.add_paragraph(text)
    doc.save(path)
    return filename

def generate_ppt(text, filename="output.pptx"):
    path = f"output/{filename}"
    prs = Presentation()
    
    # Split text into slides (each slide will hold ~300 characters)
    max_chars_per_slide = 300  
    slides = [text[i:i + max_chars_per_slide] for i in range(0, len(text), max_chars_per_slide)]

    for slide_text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))  # Position & size
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.add_paragraph()
        p.text = slide_text
        p.font.size = Pt(24)  # Larger font for better readability

    prs.save(path)
    return filename
